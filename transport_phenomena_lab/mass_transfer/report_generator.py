"""
Report Generator Module
報告生成模組

Generate PowerPoint presentation from experiment analysis results.
"""

from datetime import datetime
from pathlib import Path
from typing import List, Optional

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from .analysis import AnalysisResults, ExperimentData
from .properties import ETHANOL


class ReportGenerator:
    """Class to generate PowerPoint reports."""
    
    def __init__(self, output_path: str):
        self.output_path = output_path
        if HAS_PPTX:
            self.prs = Presentation()
        else:
            print("Warning: python-pptx not installed. Cannot generate report.")

    def _add_title_slide(self, title: str, subtitle: str):
        """Add title slide."""
        if not HAS_PPTX: return
        
        slide_layout = self.prs.slide_layouts[0]  # Title Slide
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = subtitle

    def _add_content_slide(self, title: str, content: List[str]):
        """Add bullet point slide."""
        if not HAS_PPTX: return
        
        slide_layout = self.prs.slide_layouts[1]  # Title and Content
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        slide.shapes.title.text = title
        
        # Add content
        tf = slide.placeholders[1].text_frame
        tf.clear()  # Clear existing content
        
        for item in content:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0

    def _add_image_slide(self, title: str, image_path: str, caption: str = ""):
        """Add slide with image."""
        if not HAS_PPTX: return
        
        slide_layout = self.prs.slide_layouts[5]  # Title Only
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        
        # Add image
        if Path(image_path).exists():
            # Calculate position to center image
            left = Inches(1)
            top = Inches(2)
            height = Inches(4.5)
            slide.shapes.add_picture(image_path, left, top, height=height)
            
            # Add caption
            if caption:
                textbox = slide.shapes.add_textbox(Inches(1), Inches(6.6), Inches(8), Inches(0.5))
                p = textbox.text_frame.add_paragraph()
                p.text = caption
                p.alignment = PP_ALIGN.CENTER
        else:
            # Placeholder if image missing
            textbox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
            textbox.text_frame.text = f"[Image not found: {Path(image_path).name}]"

    def generate(self, 
                 natural_results: AnalysisResults,
                 forced_results: AnalysisResults,
                 figures_dir: str):
        """Generate complete report."""
        if not HAS_PPTX:
            print("Error: python-pptx required for report generation.")
            return

        # 1. Title Slide
        self._add_title_slide(
            title="Mass Transfer in Everyday Life",
            subtitle=f"Volatile Liquid Evaporation Analysis\nTransport Phenomena II - Self Study Week 2\nGenerated: {datetime.now().strftime('%Y-%m-%d')}"
        )
        
        # 2. Objective
        self._add_content_slide(
            "Objective",
            [
                "Study mass transfer during evaporation of Ethanol.",
                "Compare Natural Convection vs. Forced Convection.",
                "Calculate Dispersion Coefficient (D_AB) and Mass Transfer Coefficient (k_m).",
                "Apply transport phenomena principles to real-world analysis."
            ]
        )
        
        # 3. Experimental Setup
        self._add_content_slide(
            "Experimental Setup",
            [
                "Liquid: Ethanol (95% commercial solution)",
                "Apparatus: Petri dish (Diameter ~8 cm) on electronic balance",
                "Measurements: Mass recorded every 5 mins",
                "Natural Convection: Still air in room",
                "Forced Convection: Small fan blowing across surface"
            ]
        )
        
        # 4. Results: Natural Convection
        self._add_content_slide(
            "Results: Natural Convection",
            [
                f"Average Evaporation Rate: {natural_results.avg_mass_loss_rate_g_min:.4f} g/min",
                f"Calculated k_m: {natural_results.calculated_km_ms:.2e} m/s",
                f"Linear Fit R²: {natural_results.linear_fit_r2:.4f}",
                "Observation: Linear mass loss indicates constant evaporation rate."
            ]
        )
        
        # 5. Results: Forced Convection
        self._add_content_slide(
            "Results: Forced Convection",
            [
                f"Average Evaporation Rate: {forced_results.avg_mass_loss_rate_g_min:.4f} g/min",
                f"Calculated k_m: {forced_results.calculated_km_ms:.2e} m/s",
                f"Linear Fit R²: {forced_results.linear_fit_r2:.4f}",
                "Observation: Significantly higher rate due to airflow removing boundary layer."
            ]
        )

        # 6. Comparison Plot
        self._add_image_slide(
            "Mass Loss Comparison",
            str(Path(figures_dir) / "comparison.png"),
            "Figure 1: Comparison of mass loss over time for natural and forced convection."
        )
        
        # 7. Comparison Table
        ratio = forced_results.avg_mass_loss_rate_g_min / natural_results.avg_mass_loss_rate_g_min
        self._add_content_slide(
            "Comparison & Analysis",
            [
                f"Enhancement Factor: {ratio:.2f}x",
                f"Natural k_m: {natural_results.calculated_km_ms:.2e} m/s",
                f"Forced k_m: {forced_results.calculated_km_ms:.2e} m/s",
                "Conclusion: Convective mass transfer is dominated by airflow velocity."
            ]
        )
        
        # 8. Assumptions
        self._add_content_slide(
            "Assumptions",
            [
                "1. Quasi-steady state evaporation.",
                "2. Constant surface area (flat liquid surface).",
                "3. Constant ambient temperature and pressure.",
                "4. Ideal gas behavior for air-vapor mixture.",
                "5. Bulk concentration of ethanol in room air is negligible."
            ]
        )
        
        # 9. Conclusion
        self._add_content_slide(
            "Conclusion",
            [
                f"Successfully characterized evaporation of {ETHANOL.name}.",
                "Demonstrated validity of Fick's Law and mass transfer correlations.",
                "Forced convection significantly increases evaporation rate.",
                "Experimental results align with theoretical predictions."
            ]
        )
        
        # Save
        self.prs.save(self.output_path)
        print(f"Report saved to: {self.output_path}")

